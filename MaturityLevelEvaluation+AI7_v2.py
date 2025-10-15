# MaturityLevelEvaluation+AI6_v11.5.py
# Full app (fixed): robust JSON parsing, normalized structure, pretty baseball cards,
# consolidated roadmap, diagrams, PPTX export, debug raw outputs saved.

import streamlit as st
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import numpy as np
import json, re, os
from io import BytesIO
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# ---- App Configuration ----------------------
st.set_page_config(page_title="Cloud & AI Maturity Evaluator", layout="wide")
st.title("Cloud & Data Maturity Evaluator")
st.markdown("Assess maturity, generate executive & technical guidance, and produce baseball-card project summaries and a consolidated roadmap.")

# ---- Set your OpenAI key --------------------
api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    st.error("⚠️ OpenAI API key not found! Please create a .env file with your OPENAI_API_KEY. See .env_template for reference.")
    st.stop()

client = OpenAI(api_key=api_key)

# -------------------- CSS --------------------
st.markdown("""
<style>
  .category-header { background: linear-gradient(90deg,#1976d2,#42a5f5); color:white; padding:6px; border-radius:6px; font-weight:700; margin-bottom:6px; }
  div.stButton > button, div.stDownloadButton > button {
    background-color: #1976d2 !important;
    color: white !important;
    border-radius: 6px !important;
    padding: 8px 14px !important;
    font-weight: 600 !important;
  }
  div.stButton > button:hover, div.stDownloadButton > button:hover {
    background-color: #1565c0 !important;
    color: white !important;
  }
</style>
""", unsafe_allow_html=True)

# -------------------- Structures --------------------
levels = {1: "Greenfield", 2: "Emerging", 3: "Developing", 4: "Established", 5: "Optimized"}
categories_structure = {
    "Cloud Architecture": ["Infrastructure Design","Scalability & Performance","Multi-cloud Strategy","Cost Optimization","Disaster Recovery","Service Architecture"],
    "Data Management": ["Data Quality","Data Integration","Master Data Management","Data Lifecycle","Data Storage Strategy","Real-time Processing"],
    "Data Visualization & Insights": ["Dashboard Design","Data Storytelling","Interactive Visualizations","Advanced Analytics Techniques","Self-Service Analytics","Insight Communication"],
    "AI/ML Integration": ["Model Development","MLOps & Deployment","AI Ethics & Bias","Business Integration","AutoML Capabilities","AI Governance"],
    "Governance & Security": ["Data Privacy","Compliance Management","Access Controls","Risk Management","Audit & Monitoring","Policy Enforcement"],
    "Business Engagement": ["Stakeholder Alignment","Change Management","Skills & Training","Value Measurement","Business Process Integration","Strategic Planning"]
}

# -------------------- Sidebar inputs --------------------
st.sidebar.header("Company Context")
industry = st.sidebar.selectbox("Industry", ["Homebuilding & Real Estate","Healthcare","Manufacturing","Financial Services","Logistics","Retail","Food and Beverage"])
company_size = st.sidebar.text_input("Company size", "1,200 employees")
it_size = st.sidebar.text_input("IT department size", "50")
uses_cloud = st.sidebar.radio("Uses cloud?", ["No","Yes"], index=1)
cloud_platform = st.sidebar.text_input("Which cloud platform(s)?", "Azure") if uses_cloud == "Yes" else ""
priority_projects = st.sidebar.text_area("Priority projects", "ERP consolidation, eCommerce upgrade")
use_seed_scenario = st.sidebar.checkbox("Seed with charitable gaming scenario", value=True)
seed_scenario_text = (
    "The client is a manufacturer and distributor of charitable gaming products. "
    "They operate three business units with silos, ~10 ERPs, no consolidated data, and many long-tenured staff resistant to change."
) if use_seed_scenario else ""

# -------------------- Sliders UI --------------------
st.markdown("---")
st.markdown("## Maturity Assessment")
st.markdown("**Scale:** 1 = Greenfield | 2 = Emerging | 3 = Developing | 4 = Established | 5 = Optimized")
all_scores, category_comments, category_inclusion = {}, {}, {}
for category, sub_caps in categories_structure.items():
    with st.expander(category, expanded=False):
        st.markdown(f'<div class="category-header">{category}</div>', unsafe_allow_html=True)
        include_cat = st.checkbox(f"Include {category}", True, key=f"include_{category}")
        category_inclusion[category] = include_cat

        sub_scores = {}
        cols = st.columns(3)
        for i, sub_cap in enumerate(sub_caps):
            with cols[i % 3]:
                score = st.slider(f"{sub_cap}", 1, 5, 3, key=f"{category}_{sub_cap}", format="Level %d")
                st.caption(f"**{levels[score]}**")
                sub_scores[sub_cap] = score
            if (i+1) % 3 == 0 and i < len(sub_caps)-1:
                cols = st.columns(3)

        all_scores[category] = {"average": round(np.mean(list(sub_scores.values())), 1), "sub_capabilities": sub_scores}
        comment = st.text_area(f"Comments for {category} (optional):", key=f"comment_{category}", height=70)
        category_comments[category] = comment

overall_input = st.text_area("Overall context/constraints (budget, compliance, culture):", height=100)

# -------------------- Session-state init --------------------
if "recommendation_data" not in st.session_state: st.session_state["recommendation_data"] = []
if "category_fragments" not in st.session_state: st.session_state["category_fragments"] = []
if "consolidated_json" not in st.session_state: st.session_state["consolidated_json"] = None
if "raw_ai_outputs" not in st.session_state: st.session_state["raw_ai_outputs"] = {}

# -------------------- Helpers: OpenAI and JSON parsing --------------------
def call_openai(prompt, max_tokens=1400, temperature=0.6):
    if client is None:
        raise RuntimeError("OpenAI client is not configured. Add OPENAI_API_KEY.")
    resp = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}],
        temperature=temperature,
        max_tokens=max_tokens
    )
    return str(resp.choices[0].message.content)

def try_load_json(text):
    """
    Robust JSON loader with several fallbacks.
    Returns a Python object (usually dict) or raises ValueError.
    """
    if text is None:
        raise ValueError("No text provided")
    t = str(text).strip()

    # If code fence present, extract inner content
    fence = re.search(r"```(?:json)?\s*([\s\S]*?)\s*```", t, re.IGNORECASE)
    if fence:
        t = fence.group(1).strip()

    # Direct try
    try:
        return json.loads(t)
    except Exception:
        pass

    # Try replace single quotes with double quotes (common model output)
    try:
        return json.loads(t.replace("'", '"'))
    except Exception:
        pass

    # Remove trailing commas before } or ]
    try:
        cleaned = re.sub(r",\s*([}\]])", r"\1", t)
        return json.loads(cleaned)
    except Exception:
        pass

    # Extract first {...} substring
    start = t.find("{"); end = t.rfind("}")
    if start != -1 and end != -1 and end > start:
        candidate = t[start:end+1]
        try:
            return json.loads(candidate)
        except Exception:
            # last-ditch: replace single quotes in candidate
            try:
                return json.loads(candidate.replace("'", '"'))
            except Exception:
                pass

    raise ValueError("Could not parse JSON from the model output.")

def normalize_baseball_card(parsed):
    """
    Ensure returned object has 'executive' and 'technical' keys.
    Accept a few common variants. Returns normalized dict:
    { "executive": {...}, "technical": {...} }
    """
    if parsed is None:
        return {"executive": {}, "technical": {}}
    if isinstance(parsed, str):
        # cannot parse — return empty and keep raw elsewhere
        return {"executive": {}, "technical": {}}
    if isinstance(parsed, list):
        # unexpected — place in executive.summary
        return {"executive": {"summary": " ".join(map(str, parsed))}, "technical": {}}
    if isinstance(parsed, dict):
        keys_lower = {k.lower(): k for k in parsed.keys()}
        # If already has exec/technical
        if "executive" in parsed and "technical" in parsed:
            return {
                "executive": parsed.get("executive") or {},
                "technical": parsed.get("technical") or {}
            }
        # Accept capitalized variants
        if "Executive" in parsed or "Technical" in parsed:
            return {
                "executive": parsed.get("Executive") or parsed.get("executive") or {},
                "technical": parsed.get("Technical") or parsed.get("technical") or {}
            }
        # Some outputs may return top-level fields for executive only
        # Heuristic: if keys include summary/recommendation/activities -> treat as executive
        exec_keys = {"summary", "recommendation", "activities", "project_activities", "focus_8w", "plan_3y", "assumptions", "team"}
        lower_keys = {k.lower() for k in parsed.keys()}
        if lower_keys & exec_keys:
            # map fields to canonical names if necessary
            exec_block = {}
            tech_block = {}
            for k, v in parsed.items():
                kl = k.lower()
                if kl in exec_keys:
                    # unify 'project_activities' -> 'activities'
                    if kl == "project_activities":
                        exec_block.setdefault("activities", v)
                    else:
                        exec_block[kl] = v
                else:
                    # put other keys under exec by default
                    exec_block[k] = v
            return {"executive": exec_block, "technical": tech_block}
        # If parsed contains exactly two top-level keys that look like cards (e.g., 'Exec' and 'Tech'), map them
        if len(parsed.keys()) <= 4:
            # attempt mapping by inspection
            exec_block = parsed.get("executive") or parsed.get("Executive") or {}
            tech_block = parsed.get("technical") or parsed.get("Technical") or {}
            return {"executive": exec_block, "technical": tech_block}
        # fallback: put entire parsed content into executive.summary as string
        return {"executive": {"summary": json.dumps(parsed)[:1000]}, "technical": {}}
    # else fallback
    return {"executive": {}, "technical": {}}

def get_field(case_insensitive_dict, *candidates):
    """
    Helper: given a dict, return first existing field among candidates (case-insensitive).
    """
    if not isinstance(case_insensitive_dict, dict):
        return None
    for cand in candidates:
        for k in case_insensitive_dict.keys():
            if k.lower() == cand.lower():
                return case_insensitive_dict[k]
    return None

# -------------------- Diagram drawing --------------------
def draw_8week_roadmap_figure(focus_dict):
    fig, ax = plt.subplots(figsize=(12, 3))
    ax.set_xlim(0, 4)
    ax.set_ylim(0, 1)
    ax.axis("off")
    for i, sprint in enumerate(["sprint1", "sprint2", "sprint3", "sprint4"]):
        x = i
        items = focus_dict.get(sprint, [])
        if isinstance(items, str):
            items = [items]
        lines = [f"• {it}" for it in items] if items else ["(no items)"]
        text = f"Sprint {i+1}\n" + "\n".join(lines)
        ax.add_patch(
            patches.FancyBboxPatch((x + 0.05, 0.05), 0.9, 0.9, boxstyle="round,pad=0.02", facecolor="#e3f2fd", edgecolor="#1976d2")
        )
        ax.text(x + 0.08, 0.5, text, ha="left", va="center", fontsize=8, wrap=True)
    plt.tight_layout()
    return fig

def draw_3year_roadmap_figure(plan_dict):
    fig, ax = plt.subplots(figsize=(12, 3))
    ax.set_xlim(0, 3)
    ax.set_ylim(0, 1)
    ax.axis("off")
    for i, year in enumerate(["year1", "year2", "year3"]):
        x = i
        items = plan_dict.get(year, [])
        if isinstance(items, str):
            items = [items]
        lines = [f"• {it}" for it in items] if items else ["(no items)"]
        text = f"Year {i+1}\n" + "\n".join(lines)
        ax.add_patch(
            patches.FancyBboxPatch((x + 0.05, 0.05), 0.9, 0.9, boxstyle="round,pad=0.02", facecolor="#e8f5e9", edgecolor="#2e7d32")
        )
        ax.text(x + 0.08, 0.5, text, ha="left", va="center", fontsize=8, wrap=True)
    plt.tight_layout()
    return fig

# -------------------- PPTX helpers --------------------
def add_wrapped_paragraph(frame, text, font_size=11, bold=False, level=0):
    p = frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.level = level
    p.word_wrap = True
    return p

def export_to_pptx(consolidated, fig1, fig2, rec_data):
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Consolidated Roadmap & Baseball Cards"

    # Roadmap slides
    for title, fig in [("8-Week Roadmap", fig1), ("3-Year Roadmap", fig2)]:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        img = BytesIO()
        fig.savefig(img, format="png", bbox_inches="tight")
        img.seek(0)
        slide.shapes.add_picture(img, Inches(0.5), Inches(1.5), width=Inches(8))

    # Per-category baseball card slides
    for item in rec_data:
        cat = item["category"]
        data = item["data_normalized"]  # normalized form we saved
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"{cat} Baseball Cards"

        # Left column: Executive
        tf = slide.shapes.add_textbox(Inches(0.3), Inches(1.3), Inches(4.2), Inches(5)).text_frame
        tf.clear()
        add_wrapped_paragraph(tf, "EXECUTIVE Baseball Card", 14, True)
        exec_block = data.get("executive", {}) or {}
        # summary + recommendation
        summary = get_field(exec_block, "summary")
        rec = get_field(exec_block, "recommendation")
        if summary: add_wrapped_paragraph(tf, f"Summary: {summary}", 11)
        if rec: add_wrapped_paragraph(tf, f"Recommendation: {rec}", 11)
        activities = get_field(exec_block, "activities", "project_activities")
        if activities and isinstance(activities, list):
            add_wrapped_paragraph(tf, "Project Activities:", 11, True)
            for a in activities:
                add_wrapped_paragraph(tf, f"• {a}", 10, False, 1)
        assumptions = get_field(exec_block, "assumptions")
        if assumptions and isinstance(assumptions, list):
            add_wrapped_paragraph(tf, "Assumptions:", 11, True)
            for a in assumptions:
                add_wrapped_paragraph(tf, f"• {a}", 10, False, 1)

        # Right column: Technical
        tf2 = slide.shapes.add_textbox(Inches(4.8), Inches(1.3), Inches(4.2), Inches(5)).text_frame
        tf2.clear()
        add_wrapped_paragraph(tf2, "TECHNICAL Baseball Card", 14, True)
        tech_block = data.get("technical", {}) or {}
        summary_t = get_field(tech_block, "summary")
        rec_t = get_field(tech_block, "recommendation")
        if summary_t: add_wrapped_paragraph(tf2, f"Summary: {summary_t}", 11)
        if rec_t: add_wrapped_paragraph(tf2, f"Recommendation: {rec_t}", 11)
        t_activities = get_field(tech_block, "activities", "project_activities")
        if t_activities and isinstance(t_activities, list):
            add_wrapped_paragraph(tf2, "Project Activities:", 11, True)
            for a in t_activities:
                add_wrapped_paragraph(tf2, f"• {a}", 10, False, 1)
        assumptions_t = get_field(tech_block, "assumptions")
        if assumptions_t and isinstance(assumptions_t, list):
            add_wrapped_paragraph(tf2, "Assumptions:", 11, True)
            for a in assumptions_t:
                add_wrapped_paragraph(tf2, f"• {a}", 10, False, 1)
        # team
        team = get_field(tech_block, "team")
        if team and isinstance(team, list):
            add_wrapped_paragraph(tf2, "Initial Team (3-6 months):", 11, True)
            for t in team:
                add_wrapped_paragraph(tf2, f"• {t}", 10, False, 1)

    # final summary slide (top 3 priorities)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Executive Summary — Top Priorities"
    tf3 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(5)).text_frame
    tf3.clear()
    add_wrapped_paragraph(tf3, "Top 3 Priorities (by impact)", 18, True)

    # derive priorities from consolidated (if present)
    items = []
    if consolidated:
        for s in ["sprint1", "sprint2", "sprint3", "sprint4"]:
            items.extend(consolidated.get("focus_8w", {}).get(s, []))
        for y in ["year1", "year2", "year3"]:
            items.extend(consolidated.get("plan_3y", {}).get(y, []))
    top3 = items[:3]
    for t in top3:
        add_wrapped_paragraph(tf3, f"• {t}", 14)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# -------------------- Generate AI-powered assessment --------------------
# Build a strong prompt template that enforces required JSON schema
generation_schema = """
You are an experienced CTO advisor. Return ONLY valid JSON that exactly follows this structure (no explanatory text, no markdown fences):

{
  "executive": {
    "summary": "2-3 sentence summary",
    "recommendation": "2+ sentence justification",
    "activities": ["Activity 1", "Activity 2", "..."],
    "focus_8w": ["Sprint1 item", "Sprint2 item", "..."],
    "plan_3y": ["Year1 item", "Year2 item", "..."],
    "assumptions": ["Assumption 1", "..."]
  },
  "technical": {
    "summary": "2-3 sentence summary",
    "recommendation": "2+ sentence technical justification",
    "activities": ["Tactic 1", "Tactic 2", "..."],
    "focus_8w": ["Sprint-level technical task", "..."],
    "plan_3y": ["Year1 technical plan", "..."],
    "assumptions": ["Assumption A", "..."],
    "team": ["Role: count", "..."]
  }
}

Make sure:
- All keys are double quoted.
- All lists are JSON arrays.
- Keep entries concise.
- Use the inputs below for context.
"""

if st.button("Generate AI-Powered Strategic Assessment"):
    if client is None:
        st.error("OpenAI not configured. Add OPENAI_API_KEY.")
    else:
        # reset storage for fresh run
        st.session_state["recommendation_data"] = []
        st.session_state["category_fragments"] = []
        st.session_state["raw_ai_outputs"] = {}
        st.session_state["consolidated_json"] = None

        # select categories: include check OR comment present -> included
        categories_to_process = [
            c for c in categories_structure.keys()
            if category_inclusion.get(c) or (category_comments.get(c, "").strip() != "")
        ]
        if not categories_to_process:
            st.info("No categories selected — check 'Include' for categories to evaluate or add a comment to include it.")
        else:
            with st.spinner("Calling AI for selected categories..."):
                for category in categories_to_process:
                    include_flag = category_inclusion.get(category, False)
                    comment_text = category_comments.get(category, "").strip()
                    scores = all_scores.get(category, {})
                    avg = scores.get("average")
                    # Build a detailed prompt for the model — uses schema defined above
                    prompt = generation_schema + f"""

Context:
Industry: {industry}
Company size: {company_size}
IT department size: {it_size}
Uses cloud: {uses_cloud} {cloud_platform}
Priority projects: {priority_projects if priority_projects else 'None'}
Category: {category}
Included flag: {'Yes' if include_flag else 'No'}
Category maturity average (if included): {avg if include_flag else 'N/A'}
Sub-capability scores: {json.dumps(scores.get('sub_capabilities', {}))}
Category comments: {comment_text if comment_text else 'None'}
Overall context: {overall_input if overall_input else 'None'}
Seed scenario: {seed_scenario_text if seed_scenario_text else 'None'}

Return the JSON only, exactly matching the schema at the top.
"""
                    try:
                        raw = call_openai(prompt, max_tokens=1000, temperature=0.4)
                        st.session_state["raw_ai_outputs"][category] = raw
                        # parse robustly
                        parsed = try_load_json(raw)
                        normalized = normalize_baseball_card(parsed)
                        # store both raw, parsed and normalized for debugging & export
                        st.session_state["recommendation_data"].append({
                            "category": category,
                            "raw": raw,
                            "parsed": parsed,
                            "data_normalized": normalized,
                            "show_avg": include_flag,
                            "avg": avg if include_flag else None
                        })
                        # For consolidation, use executive.focus_8w and plan_3y if present (normalize to list)
                        exec_focus = get_field(normalized["executive"], "focus_8w") or []
                        exec_plan3 = get_field(normalized["executive"], "plan_3y") or []
                        if isinstance(exec_focus, str): exec_focus = [exec_focus]
                        if isinstance(exec_plan3, str): exec_plan3 = [exec_plan3]
                        st.session_state["category_fragments"].append({
                            "category": category,
                            "focus_8w": exec_focus,
                            "plan_3y": exec_plan3
                        })
                    except Exception as e:
                        st.error(f"Failed to generate/parse JSON for '{category}': {e}")
                        # save raw text for debugging if available
                        st.session_state["raw_ai_outputs"][category] = raw if 'raw' in locals() else "<no raw captured>"

# -------------------- Display pretty Baseball Cards --------------------
if st.session_state.get("recommendation_data"):
    st.markdown("---")
    st.markdown("## AI-generated Baseball Cards (Executive & Technical)")
    for item in st.session_state["recommendation_data"]:
        cat = item["category"]
        normalized = item.get("data_normalized", {})
        raw_text = item.get("raw", "")
        st.subheader(cat)
        # Show maturity level when included
        if item.get("show_avg") and item.get("avg") is not None:
            level_label = levels.get(int(round(item["avg"])), "")
            st.caption(f"Reported maturity average: {item['avg']} — {level_label}")

        # EXECUTIVE card
        st.markdown("**EXECUTIVE Baseball Card**")
        exec_block = normalized.get("executive", {}) or {}
        if exec_block:
            summary = get_field(exec_block, "summary")
            recommendation = get_field(exec_block, "recommendation")
            activities = get_field(exec_block, "activities", "project_activities")
            focus8 = get_field(exec_block, "focus_8w")
            plan3 = get_field(exec_block, "plan_3y")
            assumptions = get_field(exec_block, "assumptions")

            if summary: st.markdown(f"- **Summary:** {summary}")
            if recommendation: st.markdown(f"- **Recommendation:** {recommendation}")
            if activities and isinstance(activities, list):
                st.markdown("- **Project Activities:**")
                for a in activities: st.markdown(f"  • {a}")
            if focus8 and isinstance(focus8, list):
                st.markdown("- **8-Week Focus:**")
                for f in focus8: st.markdown(f"  • {f}")
            if plan3 and isinstance(plan3, list):
                st.markdown("- **3-Year Plan:**")
                for p in plan3: st.markdown(f"  • {p}")
            if assumptions and isinstance(assumptions, list):
                st.markdown("- **Assumptions:**")
                for a in assumptions: st.markdown(f"  • {a}")
        else:
            st.info("No Executive card generated.")
            with st.expander(f"Raw AI output for '{cat}' (executive missing)"):
                st.code(raw_text)

        st.markdown("---")
        # TECHNICAL card
        st.markdown("**TECHNICAL Baseball Card**")
        tech_block = normalized.get("technical", {}) or {}
        if tech_block:
            summary = get_field(tech_block, "summary")
            recommendation = get_field(tech_block, "recommendation")
            activities = get_field(tech_block, "activities", "project_activities")
            focus8 = get_field(tech_block, "focus_8w")
            plan3 = get_field(tech_block, "plan_3y")
            assumptions = get_field(tech_block, "assumptions")
            team = get_field(tech_block, "team")

            if summary: st.markdown(f"- **Summary:** {summary}")
            if recommendation: st.markdown(f"- **Recommendation:** {recommendation}")
            if activities and isinstance(activities, list):
                st.markdown("- **Project Activities:**")
                for a in activities: st.markdown(f"  • {a}")
            if focus8 and isinstance(focus8, list):
                st.markdown("- **8-Week Tactical Plan:**")
                for f in focus8: st.markdown(f"  • {f}")
            if plan3 and isinstance(plan3, list):
                st.markdown("- **3-Year Technical Roadmap:**")
                for p in plan3: st.markdown(f"  • {p}")
            if assumptions and isinstance(assumptions, list):
                st.markdown("- **Assumptions:**")
                for a in assumptions: st.markdown(f"  • {a}")
            if team and isinstance(team, list):
                st.markdown("- **Initial Team (3–6 months):**")
                for t in team: st.markdown(f"  • {t}")
        else:
            st.info("No Technical card generated.")
            with st.expander(f"Raw AI output for '{cat}' (technical missing)"):
                st.code(raw_text)

# -------------------- Consolidate Roadmap (button) --------------------
st.markdown("---")
st.markdown("## Consolidated Roadmap")
if not st.session_state.get("category_fragments"):
    st.info("No roadmap fragments yet — generate AI recommendations first for at least one category (Include it or add a comment).")

if st.session_state.get("category_fragments"):
    if st.button("Show Consolidated Roadmap"):
        # build consolidation prompt
        fragments = st.session_state["category_fragments"]
        prompt = f"""
You are a CTO. Consolidate these category-level fragments into ONE JSON roadmap. Return ONLY JSON matching this structure:

{{
  "focus_8w": {{
    "sprint1": ["..."],
    "sprint2": ["..."],
    "sprint3": ["..."],
    "sprint4": ["..."]
  }},
  "plan_3y": {{
    "year1": ["..."],
    "year2": ["..."],
    "year3": ["..."]
  }}
}}

Category fragments:
{json.dumps(fragments, indent=2)}

Distribute initiatives sensibly across sprints and years. Return JSON only.
"""
        try:
            raw = call_openai(prompt, max_tokens=800, temperature=0.4)
            st.session_state["raw_ai_outputs"]["consolidate"] = raw
            consolidated = try_load_json(raw)
            # normalize structure
            consolidated.setdefault("focus_8w", {})
            consolidated.setdefault("plan_3y", {})
            for s in ["sprint1", "sprint2", "sprint3", "sprint4"]:
                if s not in consolidated["focus_8w"] or not isinstance(consolidated["focus_8w"][s], list):
                    consolidated["focus_8w"][s] = []
            for y in ["year1", "year2", "year3"]:
                if y not in consolidated["plan_3y"] or not isinstance(consolidated["plan_3y"][y], list):
                    consolidated["plan_3y"][y] = []
            st.session_state["consolidated_json"] = consolidated
        except Exception as e:
            st.error(f"Failed to consolidate roadmap: {e}")
            with st.expander("Raw consolidation output"):
                st.write(st.session_state["raw_ai_outputs"].get("consolidate", "<no raw>"))

# If consolidated exists in session_state, show diagrams and allow PPTX export (persist after download)
if st.session_state.get("consolidated_json"):
    consolidated = st.session_state["consolidated_json"]
    fig1 = draw_8week_roadmap_figure(consolidated.get("focus_8w", {}))
    fig2 = draw_3year_roadmap_figure(consolidated.get("plan_3y", {}))

    st.markdown("### 8-Week Roadmap Diagram")
    st.pyplot(fig1)

    st.markdown("### 3-Year Roadmap Diagram")
    st.pyplot(fig2)

    # pretty print consolidated text as well
    st.markdown("### Consolidated 8-Week Focus")
    for s in ["sprint1", "sprint2", "sprint3", "sprint4"]:
        st.markdown(f"**{s.capitalize()}**")
        for it in consolidated["focus_8w"].get(s, []):
            st.markdown(f"- {it}")

    st.markdown("### Consolidated 3-Year Plan")
    for y in ["year1", "year2", "year3"]:
        st.markdown(f"**{y.capitalize()}**")
        for it in consolidated["plan_3y"].get(y, []):
            st.markdown(f"- {it}")

    # PPTX export (cons + per-category normalized cards)
    try:
        pptx_bytes = export_to_pptx(consolidated, fig1, fig2, st.session_state["recommendation_data"])
        st.download_button("📥 Download Roadmap and Baseball Cards (PowerPoint)", data=pptx_bytes,
                           file_name="Consolidated_Roadmap_and_Cards.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    except Exception as e:
        st.error(f"PPTX export failed: {e}")

st.markdown("---")
